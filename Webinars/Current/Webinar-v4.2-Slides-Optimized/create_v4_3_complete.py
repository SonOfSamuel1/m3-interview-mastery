#!/usr/bin/env python3
"""
Create Complete Webinar v4.3 PowerPoint for Canva
- All 214 text slides from markdown
- All 27 combo slides from COMBO_SLIDES data
- No markdown syntax (clean bold formatting)
- Rich formatting applied
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re

# v4.2 Color Scheme
PRIMARY_COLOR = RGBColor(26, 54, 93)      # #1a365d
SECONDARY_COLOR = RGBColor(44, 82, 130)   # #2c5282
ACCENT_COLOR = RGBColor(49, 130, 206)     # #3182ce
TEXT_COLOR = RGBColor(26, 32, 44)         # #1a202c
GRAY_DARK = RGBColor(74, 85, 104)         # #4a5568

# Combo Slides Data (from generate_slides_v4_2.py)
COMBO_SLIDES = {
    '12-img': {
        'type': 'narrative_illustration',
        'top_text': '10 Years. $130K. Stuck.',
        'subtitle': 'AT&T - 2021',
        'image_desc': 'Before state: AT&T office environment. "$130K after 10 years" text overlay. Stressed professional at desk.'
    },
    '13B-img': {
        'type': 'narrative_illustration',
        'top_text': '90 Days From Now...',
        'bottom_text': 'You\'re Reviewing THREE Job Offers',
        'image_desc': 'Clean, organized home workspace scene. Morning light streaming in. Quality coffee mug. Laptop showing multiple opportunity tabs.'
    },
    '13D-img': {
        'type': 'narrative_illustration',
        'top_text': 'YOU Get to Choose:',
        'bullets': ['Company Culture', 'Role Fit', 'Compensation', 'Work-Life Balance'],
        'image_desc': 'Split screen showing: company culture icons, role descriptions, compensation packages, work-life balance symbols.'
    },
    '13F-img': {
        'type': 'narrative_illustration',
        'top_text': 'This Isn\'t Fantasy. This Is Real.',
        'bottom_text': 'And I\'ll Show You Exactly How',
        'image_desc': 'Before/After transformation visual. Left: stressed person at cluttered desk. Right: confident person reviewing multiple offers.'
    },
    '22-img': {
        'type': 'data_visualization',
        'headline': '10 Years of Hard Work',
        'metric': '$65K → $130K',
        'bottom_text': 'But I Hit a Ceiling...',
        'image_desc': 'Timeline graphic: 2011-2021. AT&T logo. Income progression visualization showing growth plateau.'
    },
    '23-img': {
        'type': 'data_visualization',
        'headline': '$130K = The Ceiling',
        'bottom_text': 'No Matter How Hard I Worked',
        'image_desc': 'Income ceiling graph. Growth line plateauing. Visual representation of career plateau.'
    },
    '24-img': {
        'type': 'narrative_illustration',
        'top_text': 'And It Wasn\'t Enough',
        'corners': ['Debt Piling Up', 'Missing Family Time', '60-Hour Weeks', 'Career Stuck'],
        'image_desc': 'Four-panel pain point visual: debt bills stack, missed family moments, long work hours clock, career stagnation barrier.'
    },
    '26-img': {
        'type': 'narrative_illustration',
        'top_text': 'Then I Overheard Two Coworkers Talking...',
        'bottom_text': 'Same Skills as Me. Different Company. 3× the Income.',
        'image_desc': 'Office environment. Two silhouetted colleagues. Speech bubble: "$400,000". Shocked listener reaction.'
    },
    '37-img': {
        'type': 'framework_diagram',
        'title': 'The 3 Positioning Factors',
        'bottom_text': 'Master These → Control Your Income',
        'image_desc': 'Three architectural pillars. Left: "WHO - Which Companies Pay More". Center: "WHICH - Roles & Comp Structure". Right: "HOW - Strategic Positioning".'
    },
    '46-img': {
        'type': 'narrative_illustration',
        'top_text': 'Sound Familiar?',
        'image_desc': 'Thought bubbles surrounding central figure: "Not qualified enough", "No connections", "Don\'t have the background".'
    },
    '48-img': {
        'type': 'before_after',
        'top_text': 'Here\'s What You Need to Hear...',
        'left_title': 'YOU ALREADY HAVE:',
        'left_items': ['✓ Experience', '✓ Skills', '✓ Track Record'],
        'right_title': 'YOU DON\'T NEED:',
        'right_items': ['✗ MBA', '✗ Connections', '✗ More Skills'],
        'bottom_text': 'You\'re Not Under-Qualified. You\'re Under-POSITIONED.',
    },
    '51-img': {
        'type': 'data_visualization',
        'headline': 'Week 3: Results',
        'bottom_text': '3 Interviews Scheduled',
        'image_desc': 'Success timeline. Week 1-3 progression. Google, Amazon, Salesforce logos with checkmarks.'
    },
    '56-img': {
        'type': 'narrative_illustration',
        'top_text': 'The Doubts Came Back...',
        'bottom_text': 'Have You Felt This Way?',
        'image_desc': 'Impostor syndrome visualization. Central defeated figure. Surrounding thought bubbles with doubt statements.'
    },
    '59-img': {
        'type': 'narrative_illustration',
        'top_text': 'The Decision:',
        'bottom_text': 'This Time, I\'ll Be Prepared',
        'subtitle': 'Positioning = Preparation',
        'image_desc': 'Determined professional at organized desk. Research materials spread out. Preparation checklist visible.'
    },
    '64-img': {
        'type': 'before_after',
        'title': 'YOU ARE QUALIFIED',
        'subtitle': 'The Problem Isn\'t Your Skills... It\'s Your Positioning',
    },
    '66-img': {
        'type': 'framework_diagram',
        'title': 'You Don\'t Need NEW Skills. You Need to REPOSITION Your Existing Skills.',
        'bottom_text': 'It\'s Not About HAVING More. It\'s About POSITIONING What You Have.',
        'image_desc': 'Skills translation matrix showing current skills repositioned for $300K+ roles.'
    },
    '73-img': {
        'type': 'data_visualization',
        'headline': 'The Pattern Is Clear',
        'bottom_text': 'Same Skills • Better Positioning • Repeatable Results',
        'image_desc': 'Three case studies side-by-side showing repeatable pattern.'
    },
    '75-img': {
        'type': 'narrative_illustration',
        'top_text': 'But What About...',
        'bottom_text': 'Can This Really Work for YOU?',
        'image_desc': 'External objections thought cloud. Multiple bubbles: "No time", "Full-time job", "Family responsibilities".'
    },
    '84-img': {
        'type': 'data_visualization',
        'headline': 'The 5-Day Experiment',
        'metric': '30 min/day × 5 days = 3 INTERVIEWS',
        'bottom_text': 'Same System Still Works Today',
        'image_desc': '5-day calendar with daily 30-min blocks. Leading to "Week 2: THREE INTERVIEWS SCHEDULED".'
    },
    '86-img': {
        'type': 'framework_diagram',
        'title': 'Reverse Attraction',
        'bottom_text': 'Stop Chasing. Start Attracting.',
        'image_desc': 'Process reversal diagram. Traditional vs Reverse attraction flow.'
    },
    '93-img': {
        'type': 'data_visualization',
        'headline': 'Darrell\'s First Success',
        'metric': '+$90K Raise',
        'subtitle': '$100K → $190K',
        'bottom_text': '30 Min/Day • While Employed • Weeks, Not Years',
    },
    '96-img': {
        'type': 'data_visualization',
        'headline': 'The System Works',
        'bottom_text': '30 Min/Day • Employed • $75K+ Raises',
        'image_desc': 'Three success stories comparison. Each showing proven system.'
    },
    '101-img': {
        'type': 'narrative_illustration',
        'top_text': '1 Year From Now',
        'metric': '$75K-$125K More Annually',
        'image_desc': 'One-year transformation timeline. "Today" vs "1 Year from Now" side-by-side comparison.'
    },
    '102-img': {
        'type': 'before_after',
        'left_text': 'You at 4:45 PM • Dinner with Family',
        'right_text': 'Old Colleagues • Still Working',
        'bottom_text': 'Better Position = Better Life',
    },
    '110-img': {
        'type': 'framework_diagram',
        'title': 'Total Program Value',
        'metric': '$50,500',
        'image_desc': 'Value stack calculator. All seven components listed with individual values.'
    },
    '126-img': {
        'type': 'narrative_illustration',
        'headline': '15 Spots This Month',
        'bottom_text': 'Limited Availability • Book Your Call Today',
        'image_desc': 'Scarcity visual with countdown timer or limited slot indicators.'
    },
    '139-img': {
        'type': 'data_visualization',
        'headline': 'Day 1 → Day 45',
        'bottom_text': 'Your First Interview (They Recruited YOU)',
        'subtitle': 'This Is Positioning in Action',
        'image_desc': '45-day timeline. "Day 1: Strategy Call" → "Day 45: First Interview".'
    }
}

def parse_markdown_file(filepath):
    """Parse markdown file and extract text slides"""
    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()

    pattern = r'### Slide ([\d]+[A-Z]?(?:-\d+)?)\s*(?::\s*([^\n]*))?\n(.*?)(?=### Slide [\d]+|$)'
    matches = re.findall(pattern, content, re.DOTALL)

    slides = []
    for match in matches:
        slide_num = match[0].strip()
        slide_content = match[2].strip()

        # Clean content
        lines = [line.strip() for line in slide_content.split('\n') if line.strip()]
        cleaned_content = '\n'.join(lines)
        cleaned_content = re.sub(r'\[.*?\]', '', cleaned_content)
        cleaned_content = re.sub(r'\n\s*\n', '\n', cleaned_content).strip()

        slides.append({
            'number': slide_num,
            'content': cleaned_content,
            'is_combo': False
        })

    return slides

def strip_markdown_bold(text):
    """Remove **bold** markdown syntax"""
    return re.sub(r'\*\*(.*?)\*\*', r'\1', text)

def add_formatted_text(text_frame, content, center=True):
    """Add formatted text without markdown syntax"""
    text_frame.clear()
    lines = content.split('\n') if isinstance(content, str) else [content]

    for i, line in enumerate(lines):
        line = str(line).strip()
        if not line:
            continue

        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
        p.space_after = Pt(8)

        # Check if title-like
        is_title = (line.isupper() and len(line) < 80) or 'SECRET #' in line or 'POSITIONING' in line

        # Parse markdown bold
        parts = re.split(r'(\*\*.*?\*\*)', line)

        for part in parts:
            if not part:
                continue

            was_bold = part.startswith('**') and part.endswith('**')
            clean_text = strip_markdown_bold(part)

            run = p.add_run()
            run.text = clean_text
            run.font.name = 'Helvetica Neue'

            if is_title:
                run.font.size = Pt(60)
                run.font.bold = True
                run.font.color.rgb = PRIMARY_COLOR
            elif was_bold:
                run.font.size = Pt(48)
                run.font.bold = True
                run.font.color.rgb = SECONDARY_COLOR
            elif len(clean_text) < 50:
                run.font.size = Pt(52)
                run.font.color.rgb = TEXT_COLOR
            else:
                run.font.size = Pt(40)
                run.font.color.rgb = TEXT_COLOR

def create_text_slide(prs, slide_data):
    """Create standard text slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    textbox.text_frame.word_wrap = True
    textbox.text_frame.vertical_anchor = 1

    add_formatted_text(textbox.text_frame, slide_data['content'])

def create_combo_slide(prs, slide_num, slide_data):
    """Create combo slide with text + image placeholder"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Left: Text content
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(4.5), Inches(6))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = 1

    # Build text content from slide data
    content_parts = []

    if 'headline' in slide_data:
        content_parts.append(f"**{slide_data['headline']}**")
    if 'title' in slide_data:
        content_parts.append(f"**{slide_data['title']}**")
    if 'top_text' in slide_data:
        content_parts.append(f"**{slide_data['top_text']}**")

    if 'metric' in slide_data:
        content_parts.append(f"\n{slide_data['metric']}")
    if 'subtitle' in slide_data:
        content_parts.append(slide_data['subtitle'])

    if 'bullets' in slide_data:
        for bullet in slide_data['bullets']:
            content_parts.append(f"• {bullet}")

    if 'left_title' in slide_data and 'right_title' in slide_data:
        content_parts.append(f"\n**{slide_data['left_title']}**")
        if 'left_items' in slide_data:
            for item in slide_data['left_items']:
                content_parts.append(item)
        content_parts.append(f"\n**{slide_data['right_title']}**")
        if 'right_items' in slide_data:
            for item in slide_data['right_items']:
                content_parts.append(item)

    if 'left_text' in slide_data:
        content_parts.append(f"\n{slide_data['left_text']}")
    if 'right_text' in slide_data:
        content_parts.append(f"\n{slide_data['right_text']}")

    if 'corners' in slide_data:
        for corner in slide_data['corners']:
            content_parts.append(f"• {corner}")

    if 'bottom_text' in slide_data:
        content_parts.append(f"\n{slide_data['bottom_text']}")

    text_content = '\n'.join(content_parts)
    add_formatted_text(tf, text_content, center=False)

    # Right: Image placeholder
    img_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.5), Inches(1), Inches(4), Inches(6)
    )
    img_shape.fill.background()
    img_shape.line.color.rgb = ACCENT_COLOR
    img_shape.line.width = Pt(4)
    img_shape.line.dash_style = 2

    # Add [IMAGE] label
    img_tf = img_shape.text_frame
    img_tf.clear()
    p = img_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    # Add image description
    img_desc = slide_data.get('image_desc', '[IMAGE]')
    run = p.add_run()
    run.text = f"[IMAGE]\n\n{img_desc}"
    run.font.size = Pt(24)
    run.font.color.rgb = ACCENT_COLOR
    run.font.name = 'Helvetica Neue'

def natural_sort_key(s):
    """Natural sorting for slide numbers (1, 2, 10, 12-img, 13A, 13B-img, etc.)"""
    parts = re.split(r'(\d+)', s)
    return [int(part) if part.isdigit() else part for part in parts]

def create_presentation(markdown_file, output_file):
    """Create complete PowerPoint with all slides"""
    print("Parsing markdown file...")
    text_slides = parse_markdown_file(markdown_file)

    print(f"Found {len(text_slides)} text slides")
    print(f"Found {len(COMBO_SLIDES)} combo slides")
    print(f"Total: {len(text_slides) + len(COMBO_SLIDES)} slides")

    # Create combined list with proper ordering
    all_slides = []

    # Add text slides
    for slide in text_slides:
        all_slides.append({
            'number': slide['number'],
            'type': 'text',
            'data': slide
        })

    # Add combo slides
    for slide_num, slide_data in COMBO_SLIDES.items():
        all_slides.append({
            'number': slide_num,
            'type': 'combo',
            'data': slide_data
        })

    # Sort by slide number
    all_slides.sort(key=lambda x: natural_sort_key(x['number']))

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

    print("\nCreating slides in order...")
    for i, slide_info in enumerate(all_slides, 1):
        if slide_info['type'] == 'text':
            create_text_slide(prs, slide_info['data'])
            print(f"  [{i}/{len(all_slides)}] Text: {slide_info['number']}")
        else:
            create_combo_slide(prs, slide_info['number'], slide_info['data'])
            print(f"  [{i}/{len(all_slides)}] Combo: {slide_info['number']}")

    print(f"\nSaving to {output_file}...")
    prs.save(output_file)

    text_count = len(text_slides)
    combo_count = len(COMBO_SLIDES)
    total_count = len(all_slides)

    return total_count, text_count, combo_count

if __name__ == "__main__":
    markdown_file = "../Webinar-Script-v4.2-Optimized.md"
    output_file = "Webinar-v4.3-Canva-COMPLETE.pptx"

    total, text, combo = create_presentation(markdown_file, output_file)

    print(f"""
✓ POWERPOINT CREATED SUCCESSFULLY!

File: {output_file}
Total Slides: {total}
  • Text slides: {text}
  • Combo slides: {combo}

FORMATTING APPLIED:
  ✓ NO markdown syntax (** removed, actual bold applied)
  ✓ Title text: 60pt bold in primary blue
  ✓ Bold emphasis: 48pt bold in secondary blue
  ✓ Regular text: 40-52pt based on length
  ✓ ALL combo slides included with image placeholders
  ✓ Image descriptions included in placeholders
  ✓ Proper slide ordering maintained

READY FOR CANVA IMPORT!
""")
