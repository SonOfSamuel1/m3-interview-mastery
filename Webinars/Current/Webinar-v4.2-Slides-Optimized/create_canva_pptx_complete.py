#!/usr/bin/env python3
"""
Create complete Canva-ready PowerPoint with ALL slides
- Includes all 214 text slides
- Includes all 27 combo slides
- Removes markdown syntax (** for bold)
- Applies actual rich formatting
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

# v4.2 Color Scheme
PRIMARY_COLOR = RGBColor(26, 54, 93)      # #1a365d
SECONDARY_COLOR = RGBColor(44, 82, 130)   # #2c5282
ACCENT_COLOR = RGBColor(49, 130, 206)     # #3182ce
TEXT_COLOR = RGBColor(26, 32, 44)         # #1a202c
GRAY_DARK = RGBColor(74, 85, 104)         # #4a5568

def parse_markdown_file(filepath):
    """Parse markdown file and extract all slides (text + combo)"""
    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()

    # Match both regular slides and combo slides (with -img suffix)
    pattern = r'### Slide ([\d]+[A-Z]?(?:-\d+)?(?:-img)?)\s*(?::\s*([^\n]*))?\n(.*?)(?=### Slide [\d]+|$)'
    matches = re.findall(pattern, content, re.DOTALL)

    slides = []
    for match in matches:
        slide_num = match[0].strip()
        slide_title = match[1].strip() if match[1] else ""
        slide_content = match[2].strip()

        # Clean content
        lines = [line.strip() for line in slide_content.split('\n') if line.strip()]
        cleaned_content = '\n'.join(lines)

        # Remove delivery notes
        cleaned_content = re.sub(r'\[.*?\]', '', cleaned_content)
        cleaned_content = re.sub(r'\n\s*\n', '\n', cleaned_content).strip()

        is_combo = slide_num.endswith('-img')

        slides.append({
            'number': slide_num,
            'title': slide_title,
            'content': cleaned_content,
            'is_combo': is_combo
        })

    return slides

def strip_markdown_bold(text):
    """Remove **bold** markdown syntax and return clean text"""
    return re.sub(r'\*\*(.*?)\*\*', r'\1', text)

def is_bold_text(original_text):
    """Check if text was marked as bold in markdown"""
    return '**' in original_text

def apply_rich_formatting(text_frame, content):
    """Apply rich formatting without markdown syntax"""
    text_frame.clear()
    lines = content.split('\n')

    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue

        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(12)

        # Check if line should be emphasized (title-like)
        is_title = (line.isupper() and len(line) < 80) or \
                   'SECRET #' in line or \
                   'POSITIONING' in line or \
                   line.startswith('P.R.I.N.T.')

        # Parse markdown bold markers
        parts = re.split(r'(\*\*.*?\*\*)', line)

        for part in parts:
            if not part:
                continue

            # Check if this part was marked as bold
            was_bold = part.startswith('**') and part.endswith('**')
            clean_text = strip_markdown_bold(part)

            run = p.add_run()
            run.text = clean_text
            run.font.name = 'Helvetica Neue'

            # Apply formatting based on context
            if is_title:
                run.font.size = Pt(64)
                run.font.bold = True
                run.font.color.rgb = PRIMARY_COLOR
            elif was_bold:
                run.font.size = Pt(52)
                run.font.bold = True
                run.font.color.rgb = SECONDARY_COLOR
            elif len(clean_text) < 50:
                run.font.size = Pt(56)
                run.font.color.rgb = TEXT_COLOR
            else:
                run.font.size = Pt(42)
                run.font.color.rgb = TEXT_COLOR

def create_text_slide(prs, slide_data):
    """Create a standard text slide"""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Add text box
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = 1  # Middle

    apply_rich_formatting(text_frame, slide_data['content'])

def create_combo_slide(prs, slide_data):
    """Create a combo slide with text on one side and image placeholder on other"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Text box (left side)
    text_left = Inches(0.5)
    text_top = Inches(1.5)
    text_width = Inches(4)
    text_height = Inches(5.5)

    textbox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = 1

    apply_rich_formatting(text_frame, slide_data['content'])

    # Image placeholder box (right side)
    img_left = Inches(5.5)
    img_top = Inches(1.5)
    img_width = Inches(4)
    img_height = Inches(5.5)

    img_placeholder = slide.shapes.add_shape(
        1,  # Rectangle
        img_left, img_top, img_width, img_height
    )

    # Style the placeholder
    img_placeholder.fill.background()
    img_placeholder.line.color.rgb = ACCENT_COLOR
    img_placeholder.line.width = Pt(3)
    img_placeholder.line.dash_style = 2  # Dashed

    # Add "[IMAGE]" text in placeholder
    img_text_frame = img_placeholder.text_frame
    img_text_frame.clear()
    p = img_text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "[IMAGE]"
    p.runs[0].font.size = Pt(48)
    p.runs[0].font.color.rgb = ACCENT_COLOR
    p.runs[0].font.name = 'Helvetica Neue'

def create_presentation(markdown_file, output_file):
    """Create complete PowerPoint with all slides"""
    print("Parsing markdown file...")
    slides = parse_markdown_file(markdown_file)

    print(f"Found {len(slides)} total slides")
    text_slides = [s for s in slides if not s['is_combo']]
    combo_slides = [s for s in slides if s['is_combo']]
    print(f"  - {len(text_slides)} text slides")
    print(f"  - {len(combo_slides)} combo slides")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("\nCreating slides...")
    for i, slide_data in enumerate(slides, 1):
        if slide_data['is_combo']:
            create_combo_slide(prs, slide_data)
            print(f"  [{i}/{len(slides)}] Combo slide {slide_data['number']}")
        else:
            create_text_slide(prs, slide_data)
            print(f"  [{i}/{len(slides)}] Text slide {slide_data['number']}")

    # Save
    print(f"\nSaving to {output_file}...")
    prs.save(output_file)
    print(f"✓ Complete! {len(slides)} slides created")

    return len(slides), len(text_slides), len(combo_slides)

if __name__ == "__main__":
    markdown_file = "../Webinar-Script-v4.2-Optimized.md"
    output_file = "Webinar-v4.3-Canva-COMPLETE.pptx"

    total, text, combo = create_presentation(markdown_file, output_file)

    print(f"""
POWERPOINT CREATED SUCCESSFULLY!

File: {output_file}
Total Slides: {total}
  - Text slides: {text}
  - Combo slides: {combo}

FORMATTING APPLIED:
  ✓ No markdown syntax (** removed, actual bold applied)
  ✓ Title text: 64pt bold in primary blue
  ✓ Bold emphasis: 52pt bold in secondary blue
  ✓ Regular text: 42-56pt based on length
  ✓ All combo slides included with image placeholders
  ✓ Ready for Canva import

READY FOR CANVA!
""")
